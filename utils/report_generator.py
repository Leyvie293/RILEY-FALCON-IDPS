# utils/report_generator.py
# COMPLETE FIXED VERSION - Gets real data from both live monitors and database

import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns
from datetime import datetime, timedelta
import csv
import os
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import json
import sys
import re

sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from models import Alert, NetworkTraffic, HostActivity, db, User
from flask import current_app


class PDF(FPDF):
    def __init__(self, logo_path=None):
        super().__init__()
        self.logo_path = logo_path or os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
            'static', 'images', 'logo.png'
        )
    
    def header(self):
        if os.path.exists(self.logo_path):
            try:
                self.image(self.logo_path, x=10, y=8, w=25)
            except:
                pass
        self.set_font('Arial', 'B', 15)
        self.cell(80)
        self.cell(30, 10, 'Riley Falcon Security Services', 0, 0, 'C')
        self.ln(20)

    def footer(self):
        self.set_y(-15)
        self.set_font('Arial', 'I', 8)
        self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')


class ReportGenerator:
    """Main Report Generator class for IDPS - FIXED VERSION"""
    
    def __init__(self):
        self.report_dir = 'reports/'
        os.makedirs(self.report_dir, exist_ok=True)
        self.logo_path = os.path.join(
            os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
            'static', 'images', 'logo.png'
        )
    
    def get_data(self, start_date, end_date, data_source='database'):
        """Get data from appropriate source"""
        print(f"🔍 Getting data from {data_source} source between {start_date} and {end_date}")
        
        if data_source == 'live':
            return self.get_live_data(start_date, end_date)
        else:
            return self.get_database_data(start_date, end_date)
    
    def get_live_data(self, start_date, end_date):
        """Get live data from active monitors - FIXED VERSION WITH BETTER HOST DATA"""
        data = {
            'alerts': [],
            'network': [],
            'host': []
        }
        
        # Convert dates to naive datetime for comparison (since live data uses naive)
        if start_date and hasattr(start_date, 'tzinfo') and start_date.tzinfo:
            start_date_naive = start_date.replace(tzinfo=None)
        else:
            start_date_naive = start_date
            
        if end_date and hasattr(end_date, 'tzinfo') and end_date.tzinfo:
            end_date_naive = end_date.replace(tzinfo=None)
        else:
            end_date_naive = end_date
        
        # Get live network data from global monitor
        try:
            from utils.network_monitor import network_monitor
            if network_monitor and hasattr(network_monitor, 'captured_packets'):
                print(f"📡 Found {len(network_monitor.captured_packets)} live packets in memory")
                
                for packet in list(network_monitor.captured_packets):
                    packet_time = packet.get('timestamp', datetime.now())
                    if isinstance(packet_time, str):
                        try:
                            packet_time = datetime.fromisoformat(packet_time.replace('Z', '+00:00'))
                        except:
                            packet_time = datetime.now()
                    
                    # Handle timezone
                    if hasattr(packet_time, 'tzinfo') and packet_time.tzinfo:
                        packet_time = packet_time.replace(tzinfo=None)
                    
                    if start_date_naive <= packet_time <= end_date_naive:
                        # Convert to model-like object for consistency
                        class PacketObj:
                            pass
                        
                        pobj = PacketObj()
                        pobj.timestamp = packet_time
                        pobj.source_ip = packet.get('src_ip', '0.0.0.0')
                        pobj.destination_ip = packet.get('dst_ip', '0.0.0.0')
                        pobj.protocol = packet.get('protocol', 'TCP')
                        pobj.source_port = packet.get('src_port', 0)
                        pobj.destination_port = packet.get('dst_port', 0)
                        pobj.packet_size = packet.get('size', 0)
                        pobj.flags = packet.get('flags', '')
                        pobj.is_suspicious = packet.get('suspicious', False)
                        pobj.anomaly_score = packet.get('score', 0)
                        
                        data['network'].append(pobj)
                
                print(f"📊 Filtered to {len(data['network'])} live network packets in date range")
        except Exception as e:
            print(f"⚠️ Error getting live network data: {e}")
            import traceback
            traceback.print_exc()
        
        # FIXED: Get live host data with better collection
        try:
            from utils.host_monitor import host_monitor
            if host_monitor:
                process_count = 0
                suspicious_count = 0
                
                # Method 1: Get from process_history (historical data)
                if hasattr(host_monitor, 'process_history') and len(host_monitor.process_history) > 0:
                    print(f"🖥️ Found {len(host_monitor.process_history)} items in process_history")
                    
                    for item in list(host_monitor.process_history):
                        if isinstance(item, dict):
                            # Check if this is a process record
                            item_type = item.get('type', '')
                            
                            if item_type == 'process' or 'process_name' in item:
                                proc_time = item.get('timestamp', datetime.now())
                                
                                if hasattr(proc_time, 'tzinfo') and proc_time.tzinfo:
                                    proc_time = proc_time.replace(tzinfo=None)
                                
                                if start_date_naive <= proc_time <= end_date_naive:
                                    class HostObj:
                                        pass
                                    
                                    hobj = HostObj()
                                    hobj.timestamp = proc_time
                                    hobj.process_name = item.get('process_name', item.get('name', 'Unknown'))
                                    hobj.process_id = item.get('process_id', item.get('pid', 0))
                                    hobj.user = item.get('user', 'Unknown')
                                    hobj.cpu_usage = item.get('cpu_usage', item.get('cpu', 0))
                                    hobj.memory_usage = item.get('memory_usage', item.get('memory', 0))
                                    hobj.is_suspicious = item.get('is_suspicious', False)
                                    hobj.network_connections = item.get('network_connections', '')
                                    hobj.open_files = item.get('open_files', 0)
                                    
                                    data['host'].append(hobj)
                                    process_count += 1
                                    if hobj.is_suspicious:
                                        suspicious_count += 1
                    
                    print(f"🖥️ Found {process_count} host processes in history ({suspicious_count} suspicious) for date range")
                
                # Method 2: If no historical data, get current processes
                if process_count == 0 and hasattr(host_monitor, 'get_top_processes'):
                    print("🖥️ No historical process data, getting current processes...")
                    current_processes = host_monitor.get_top_processes(limit=200)
                    
                    for proc in current_processes:
                        class HostObj:
                            pass
                        
                        hobj = HostObj()
                        hobj.timestamp = datetime.utcnow()
                        hobj.process_name = proc.get('process_name', 'Unknown')
                        hobj.process_id = proc.get('process_id', 0)
                        hobj.user = proc.get('user', 'Unknown')
                        hobj.cpu_usage = proc.get('cpu_usage', 0)
                        hobj.memory_usage = proc.get('memory_usage', 0)
                        hobj.is_suspicious = proc.get('is_suspicious', False)
                        hobj.network_connections = ''
                        hobj.open_files = 0
                        
                        # Only include if within date range (current time should be)
                        current_time_naive = datetime.utcnow()
                        if start_date_naive <= current_time_naive <= end_date_naive:
                            data['host'].append(hobj)
                            process_count += 1
                            if hobj.is_suspicious:
                                suspicious_count += 1
                    
                    print(f"🖥️ Added {process_count} current processes to report ({suspicious_count} suspicious)")
                
                # Method 3: Try to get from database as fallback
                if process_count == 0:
                    print("🖥️ No live process data, attempting database fallback...")
                    # Query database for recent host activities
                    try:
                        recent_host = HostActivity.query.filter(
                            HostActivity.timestamp >= start_date,
                            HostActivity.timestamp <= end_date
                        ).order_by(HostActivity.timestamp.desc()).limit(500).all()
                        
                        for activity in recent_host:
                            class HostObj:
                                pass
                            
                            hobj = HostObj()
                            hobj.timestamp = activity.timestamp
                            hobj.process_name = activity.process_name
                            hobj.process_id = activity.process_id
                            hobj.user = activity.user
                            hobj.cpu_usage = activity.cpu_usage
                            hobj.memory_usage = activity.memory_usage
                            hobj.is_suspicious = activity.is_suspicious
                            hobj.network_connections = activity.network_connections
                            hobj.open_files = activity.open_files
                            
                            data['host'].append(hobj)
                            process_count += 1
                            if hobj.is_suspicious:
                                suspicious_count += 1
                        
                        print(f"🖥️ Database fallback provided {process_count} host processes ({suspicious_count} suspicious)")
                    except Exception as e:
                        print(f"⚠️ Database fallback error: {e}")
                
        except Exception as e:
            print(f"⚠️ Error getting live host data: {e}")
            import traceback
            traceback.print_exc()
        
        # Get alerts from database (most consistent source)
        try:
            alerts = Alert.query.filter(
                Alert.timestamp >= start_date,
                Alert.timestamp <= end_date
            ).order_by(Alert.timestamp.desc()).limit(200).all()
            
            data['alerts'] = alerts
            print(f"🚨 Found {len(data['alerts'])} alerts in date range")
        except Exception as e:
            print(f"⚠️ Error getting alerts: {e}")
        
        # Final summary
        print(f"📊 Live data collected: {len(data['alerts'])} alerts, {len(data['network'])} packets, {len(data['host'])} host processes")
        return data
    
    def get_database_data(self, start_date, end_date):
        """Get data from database - PRIMARY SOURCE"""
        print(f"💾 Querying database for data between {start_date} and {end_date}")
        
        try:
            # Get network traffic from database
            network_query = NetworkTraffic.query.filter(
                NetworkTraffic.timestamp >= start_date,
                NetworkTraffic.timestamp <= end_date
            ).order_by(NetworkTraffic.timestamp.desc())
            
            network = network_query.all()
            network_count = network_query.count()
            
            # Get host activity from database
            host_query = HostActivity.query.filter(
                HostActivity.timestamp >= start_date,
                HostActivity.timestamp <= end_date
            ).order_by(HostActivity.timestamp.desc())
            
            host = host_query.all()
            host_count = host_query.count()
            
            # Get alerts from database
            alert_query = Alert.query.filter(
                Alert.timestamp >= start_date,
                Alert.timestamp <= end_date
            ).order_by(Alert.timestamp.desc())
            
            alerts = alert_query.all()
            alert_count = alert_query.count()
            
            print(f"✅ Database query results:")
            print(f"   - {alert_count} alerts")
            print(f"   - {network_count} network packets")
            print(f"   - {host_count} host activities")
            
            # Get additional stats for verification
            if alert_count == 0:
                # Check if there are ANY alerts in the database
                total_alerts = Alert.query.count()
                if total_alerts > 0:
                    print(f"⚠️ No alerts in date range, but {total_alerts} total alerts exist in database")
                    # Get the latest alert date for debugging
                    latest = Alert.query.order_by(Alert.timestamp.desc()).first()
                    if latest:
                        print(f"   Latest alert timestamp: {latest.timestamp}")
            
            # FIXED: Also check if host data exists
            if host_count == 0:
                total_host = HostActivity.query.count()
                if total_host > 0:
                    print(f"⚠️ No host activities in date range, but {total_host} total host records exist")
                    latest_host = HostActivity.query.order_by(HostActivity.timestamp.desc()).first()
                    if latest_host:
                        print(f"   Latest host timestamp: {latest_host.timestamp}")
            
            return {
                'alerts': alerts,
                'network': network,
                'host': host
            }
            
        except Exception as e:
            print(f"❌ Error getting database data: {e}")
            import traceback
            traceback.print_exc()
            return {
                'alerts': [],
                'network': [],
                'host': []
            }
    
    def get_statistics(self, data):
        """Calculate comprehensive statistics from data"""
        alerts = data.get('alerts', [])
        network = data.get('network', [])
        host = data.get('host', [])
        
        print(f"📊 Calculating statistics from: {len(alerts)} alerts, {len(network)} packets, {len(host)} processes")
        
        stats = {
            'alerts': {
                'total': len(alerts),
                'critical': 0,
                'high': 0,
                'medium': 0,
                'low': 0,
                'active': 0,
                'resolved': 0,
                'by_type': {},
                'by_source': {}
            },
            'network': {
                'total_packets': len(network),
                'suspicious_packets': 0,
                'tcp': 0,
                'udp': 0,
                'icmp': 0,
                'http': 0,
                'https': 0,
                'dns': 0,
                'other': 0,
                'total_bytes': 0,
                'top_sources': {},
                'top_destinations': {}
            },
            'host': {
                'total_processes': len(host),
                'suspicious_processes': 0,
                'unique_users': set(),
                'avg_cpu': 0,
                'avg_memory': 0,
                'processes_by_user': {}
            }
        }
        
        # Process alerts
        for alert in alerts:
            severity = getattr(alert, 'severity', 'low')
            if severity == 'critical':
                stats['alerts']['critical'] += 1
            elif severity == 'high':
                stats['alerts']['high'] += 1
            elif severity == 'medium':
                stats['alerts']['medium'] += 1
            else:
                stats['alerts']['low'] += 1
            
            is_resolved = getattr(alert, 'is_resolved', False)
            if is_resolved:
                stats['alerts']['resolved'] += 1
            else:
                stats['alerts']['active'] += 1
            
            alert_type = getattr(alert, 'alert_type', 'unknown')
            source = getattr(alert, 'source', 'unknown')
            stats['alerts']['by_type'][alert_type] = stats['alerts']['by_type'].get(alert_type, 0) + 1
            stats['alerts']['by_source'][source] = stats['alerts']['by_source'].get(source, 0) + 1
        
        # Process network data
        for packet in network:
            # Check if suspicious
            is_suspicious = getattr(packet, 'is_suspicious', False)
            if is_suspicious:
                stats['network']['suspicious_packets'] += 1
            
            # Protocol classification
            protocol = getattr(packet, 'protocol', '').upper()
            dest_port = getattr(packet, 'destination_port', 0)
            
            if protocol == 'TCP':
                stats['network']['tcp'] += 1
                # Check HTTP/HTTPS
                if dest_port in [80, 8080, 8000]:
                    stats['network']['http'] += 1
                elif dest_port in [443, 8443]:
                    stats['network']['https'] += 1
            elif protocol == 'UDP':
                stats['network']['udp'] += 1
                if dest_port == 53:
                    stats['network']['dns'] += 1
            elif protocol == 'ICMP':
                stats['network']['icmp'] += 1
            else:
                stats['network']['other'] += 1
            
            # Accumulate bytes
            stats['network']['total_bytes'] += getattr(packet, 'packet_size', 0)
            
            # Track top IPs
            src = getattr(packet, 'source_ip', 'Unknown')
            dst = getattr(packet, 'destination_ip', 'Unknown')
            stats['network']['top_sources'][src] = stats['network']['top_sources'].get(src, 0) + 1
            stats['network']['top_destinations'][dst] = stats['network']['top_destinations'].get(dst, 0) + 1
        
        # Sort top IPs
        stats['network']['top_sources'] = dict(sorted(stats['network']['top_sources'].items(), 
                                                      key=lambda x: x[1], reverse=True)[:10])
        stats['network']['top_destinations'] = dict(sorted(stats['network']['top_destinations'].items(), 
                                                           key=lambda x: x[1], reverse=True)[:10])
        
        # Process host data
        total_cpu = 0
        total_memory = 0
        
        for proc in host:
            is_suspicious = getattr(proc, 'is_suspicious', False)
            if is_suspicious:
                stats['host']['suspicious_processes'] += 1
            
            user = getattr(proc, 'user', 'Unknown')
            if user and user != 'Unknown':
                stats['host']['unique_users'].add(user)
                stats['host']['processes_by_user'][user] = stats['host']['processes_by_user'].get(user, 0) + 1
            
            cpu = getattr(proc, 'cpu_usage', 0) or 0
            memory = getattr(proc, 'memory_usage', 0) or 0
            total_cpu += cpu
            total_memory += memory
        
        if stats['host']['total_processes'] > 0:
            stats['host']['avg_cpu'] = total_cpu / stats['host']['total_processes']
            stats['host']['avg_memory'] = total_memory / stats['host']['total_processes']
        
        stats['host']['unique_users'] = len(stats['host']['unique_users'])
        
        # Print summary for debugging
        print(f"📊 Statistics calculated:")
        print(f"   - Total Alerts: {stats['alerts']['total']}")
        print(f"   - Critical: {stats['alerts']['critical']}")
        print(f"   - High: {stats['alerts']['high']}")
        print(f"   - Total Packets: {stats['network']['total_packets']}")
        print(f"   - Suspicious Packets: {stats['network']['suspicious_packets']}")
        print(f"   - Total Processes: {stats['host']['total_processes']}")
        print(f"   - Suspicious Processes: {stats['host']['suspicious_processes']}")
        
        return stats
    
    def generate_csv_report(self, report_type, start_date, end_date, data_source='database'):
        """Generate CSV report with comprehensive data"""
        print(f"📊 Generating {report_type} CSV report from {start_date} to {end_date} (source: {data_source})")
        
        data = self.get_data(start_date, end_date, data_source)
        stats = self.get_statistics(data)
        
        filename = f"report_{report_type}_{start_date.strftime('%Y%m%d_%H%M%S')}.csv"
        filepath = os.path.join(self.report_dir, filename)
        
        with open(filepath, 'w', newline='', encoding='utf-8-sig') as f:
            writer = csv.writer(f)
            
            # Header
            writer.writerow(['RILEY FALCON SECURITY SERVICES - IDPS REPORT'])
            writer.writerow([f'Report Type: {report_type}'])
            writer.writerow([f'Data Source: {data_source.upper()}'])
            writer.writerow([f'Period: {start_date.strftime("%Y-%m-%d %H:%M:%S")} to {end_date.strftime("%Y-%m-%d %H:%M:%S")}'])
            writer.writerow([f'Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")} (EAT)'])
            writer.writerow([])
            
            # Executive Summary
            writer.writerow(['EXECUTIVE SUMMARY'])
            writer.writerow(['-' * 80])
            writer.writerow([f'Total Alerts: {stats["alerts"]["total"]}'])
            writer.writerow([f'  - Critical: {stats["alerts"]["critical"]}'])
            writer.writerow([f'  - High: {stats["alerts"]["high"]}'])
            writer.writerow([f'  - Medium: {stats["alerts"]["medium"]}'])
            writer.writerow([f'  - Low: {stats["alerts"]["low"]}'])
            writer.writerow([f'  - Active: {stats["alerts"]["active"]}'])
            writer.writerow([f'  - Resolved: {stats["alerts"]["resolved"]}'])
            writer.writerow([])
            writer.writerow([f'Network Traffic: {stats["network"]["total_packets"]:,} packets'])
            writer.writerow([f'  - Suspicious: {stats["network"]["suspicious_packets"]:,}'])
            writer.writerow([f'  - TCP: {stats["network"]["tcp"]:,}'])
            writer.writerow([f'  - UDP: {stats["network"]["udp"]:,}'])
            writer.writerow([f'  - ICMP: {stats["network"]["icmp"]:,}'])
            writer.writerow([f'  - HTTP: {stats["network"]["http"]:,}'])
            writer.writerow([f'  - HTTPS: {stats["network"]["https"]:,}'])
            writer.writerow([f'  - DNS: {stats["network"]["dns"]:,}'])
            writer.writerow([f'  - Data Transferred: {stats["network"]["total_bytes"] / (1024*1024):.2f} MB'])
            writer.writerow([])
            writer.writerow([f'Host Activity: {stats["host"]["total_processes"]:,} processes'])
            writer.writerow([f'  - Suspicious: {stats["host"]["suspicious_processes"]:,}'])
            writer.writerow([f'  - Unique Users: {stats["host"]["unique_users"]}'])
            writer.writerow([f'  - Avg CPU: {stats["host"]["avg_cpu"]:.1f}%'])
            writer.writerow([f'  - Avg Memory: {stats["host"]["avg_memory"]:.1f}%'])
            writer.writerow([])
            
            # Alert Types Breakdown
            if stats['alerts']['by_type']:
                writer.writerow(['ALERT TYPES BREAKDOWN'])
                writer.writerow(['-' * 80])
                writer.writerow(['Alert Type', 'Count'])
                for alert_type, count in sorted(stats['alerts']['by_type'].items(), key=lambda x: x[1], reverse=True):
                    writer.writerow([alert_type, count])
                writer.writerow([])
            
            # Alert Sources Breakdown
            if stats['alerts']['by_source']:
                writer.writerow(['ALERT SOURCES BREAKDOWN'])
                writer.writerow(['-' * 80])
                writer.writerow(['Source', 'Count'])
                for source, count in sorted(stats['alerts']['by_source'].items(), key=lambda x: x[1], reverse=True):
                    writer.writerow([source, count])
                writer.writerow([])
            
            # Top Source IPs
            if stats['network']['top_sources']:
                writer.writerow(['TOP SOURCE IP ADDRESSES'])
                writer.writerow(['-' * 80])
                writer.writerow(['IP Address', 'Packet Count'])
                for ip, count in stats['network']['top_sources'].items():
                    writer.writerow([ip, count])
                writer.writerow([])
            
            # Top Destination IPs
            if stats['network']['top_destinations']:
                writer.writerow(['TOP DESTINATION IP ADDRESSES'])
                writer.writerow(['-' * 80])
                writer.writerow(['IP Address', 'Packet Count'])
                for ip, count in stats['network']['top_destinations'].items():
                    writer.writerow([ip, count])
                writer.writerow([])
            
            # Detailed Network Traffic
            if data['network']:
                writer.writerow(['NETWORK TRAFFIC DETAILS'])
                writer.writerow(['-' * 80])
                writer.writerow(['Timestamp', 'Source IP', 'Source Port', 'Destination IP', 
                               'Destination Port', 'Protocol', 'Size', 'Suspicious', 'Score'])
                
                for packet in data['network'][:1000]:
                    timestamp = getattr(packet, 'timestamp', datetime.now())
                    if timestamp:
                        timestamp_str = timestamp.strftime('%Y-%m-%d %H:%M:%S')
                    else:
                        timestamp_str = str(datetime.now())
                    
                    score = getattr(packet, 'anomaly_score', 0)
                    
                    writer.writerow([
                        timestamp_str,
                        getattr(packet, 'source_ip', ''),
                        getattr(packet, 'source_port', ''),
                        getattr(packet, 'destination_ip', ''),
                        getattr(packet, 'destination_port', ''),
                        getattr(packet, 'protocol', ''),
                        getattr(packet, 'packet_size', 0),
                        'Yes' if getattr(packet, 'is_suspicious', False) else 'No',
                        f"{score:.3f}"
                    ])
                writer.writerow([])
            
            # Detailed Alerts
            if data['alerts']:
                writer.writerow(['ALERTS DETAILS'])
                writer.writerow(['-' * 80])
                writer.writerow(['Timestamp', 'Type', 'Severity', 'Source', 'Description', 'Resolved'])
                
                for alert in data['alerts'][:500]:
                    timestamp = getattr(alert, 'timestamp', datetime.now())
                    timestamp_str = timestamp.strftime('%Y-%m-%d %H:%M:%S') if timestamp else str(datetime.now())
                    
                    writer.writerow([
                        timestamp_str,
                        getattr(alert, 'alert_type', ''),
                        getattr(alert, 'severity', ''),
                        getattr(alert, 'source', ''),
                        getattr(alert, 'description', '')[:100],
                        'Yes' if getattr(alert, 'is_resolved', False) else 'No'
                    ])
                writer.writerow([])
            
            # Detailed Host Activities
            if data['host']:
                writer.writerow(['HOST ACTIVITY DETAILS'])
                writer.writerow(['-' * 80])
                writer.writerow(['Timestamp', 'Process Name', 'PID', 'User', 'CPU %', 'Memory %', 'Suspicious'])
                
                for proc in data['host'][:500]:
                    timestamp = getattr(proc, 'timestamp', datetime.now())
                    timestamp_str = timestamp.strftime('%Y-%m-%d %H:%M:%S') if timestamp else str(datetime.now())
                    
                    writer.writerow([
                        timestamp_str,
                        getattr(proc, 'process_name', ''),
                        getattr(proc, 'process_id', ''),
                        getattr(proc, 'user', ''),
                        f"{getattr(proc, 'cpu_usage', 0):.1f}",
                        f"{getattr(proc, 'memory_usage', 0):.1f}",
                        'Yes' if getattr(proc, 'is_suspicious', False) else 'No'
                    ])
        
        print(f"✅ CSV report saved to {filepath}")
        return filepath
    
    def generate_pdf_report(self, report_type, start_date, end_date, data_source='database'):
        """Generate PDF report"""
        print(f"📊 Generating {report_type} PDF report from {start_date} to {end_date} (source: {data_source})")
        
        data = self.get_data(start_date, end_date, data_source)
        stats = self.get_statistics(data)
        
        filename = f"report_{report_type}_{start_date.strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(self.report_dir, filename)
        
        pdf = PDF(logo_path=self.logo_path)
        pdf.add_page()
        
        # Title
        pdf.set_font('Arial', 'B', 20)
        pdf.cell(0, 20, 'Riley Falcon Security Services', 0, 1, 'C')
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, 'Intrusion Detection & Prevention System', 0, 1, 'C')
        pdf.set_font('Arial', 'B', 14)
        pdf.cell(0, 10, 'Security Report', 0, 1, 'C')
        
        # Date range
        pdf.set_font('Arial', '', 12)
        pdf.cell(0, 10, f'Period: {start_date.strftime("%Y-%m-%d")} to {end_date.strftime("%Y-%m-%d")}', 0, 1, 'C')
        pdf.cell(0, 10, f'Data Source: {data_source.upper()}', 0, 1, 'C')
        pdf.cell(0, 10, f'Generated: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")} (EAT)', 0, 1, 'C')
        
        pdf.line(10, pdf.get_y(), 200, pdf.get_y())
        pdf.ln(10)
        
        # Executive Summary
        pdf.set_font('Arial', 'B', 14)
        pdf.cell(0, 10, 'Executive Summary', 0, 1, 'L')
        pdf.set_font('Arial', '', 11)
        
        # Add host data note if present
        host_note = ""
        if stats['host']['total_processes'] > 0:
            host_note = f"\nHost Activity: {stats['host']['total_processes']:,} processes\n- Suspicious: {stats['host']['suspicious_processes']:,}\n- Unique Users: {stats['host']['unique_users']:,}\n- Average CPU: {stats['host']['avg_cpu']:.1f}%\n- Average Memory: {stats['host']['avg_memory']:.1f}%"
        
        summary_text = f"""
Total Alerts: {stats['alerts']['total']:,}
- Critical: {stats['alerts']['critical']:,}
- High: {stats['alerts']['high']:,}
- Medium: {stats['alerts']['medium']:,}
- Low: {stats['alerts']['low']:,}
- Active: {stats['alerts']['active']:,}
- Resolved: {stats['alerts']['resolved']:,}

Network Traffic: {stats['network']['total_packets']:,} packets
- Suspicious: {stats['network']['suspicious_packets']:,}
- TCP: {stats['network']['tcp']:,}
- UDP: {stats['network']['udp']:,}
- ICMP: {stats['network']['icmp']:,}
- HTTP: {stats['network']['http']:,}
- HTTPS: {stats['network']['https']:,}
- DNS: {stats['network']['dns']:,}
- Data Transferred: {stats['network']['total_bytes'] / (1024*1024):.2f} MB
{host_note}
        """
        
        pdf.multi_cell(0, 10, summary_text)
        
        # Network Traffic Details
        if data['network'] and len(data['network']) > 0:
            pdf.add_page()
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(0, 10, f'Network Traffic Details (Top {min(50, len(data["network"]))})', 0, 1, 'L')
            pdf.set_font('Arial', '', 9)
            
            # Table header
            pdf.set_font('Arial', 'B', 9)
            pdf.cell(30, 8, 'Time', 1)
            pdf.cell(35, 8, 'Source IP', 1)
            pdf.cell(35, 8, 'Dest IP', 1)
            pdf.cell(20, 8, 'Protocol', 1)
            pdf.cell(20, 8, 'Size', 1)
            pdf.cell(20, 8, 'Status', 1)
            pdf.ln()
            
            pdf.set_font('Arial', '', 8)
            for packet in data['network'][:50]:
                timestamp = getattr(packet, 'timestamp', datetime.now())
                timestamp_str = timestamp.strftime('%H:%M:%S') if timestamp else ''
                
                pdf.cell(30, 6, timestamp_str, 1)
                pdf.cell(35, 6, getattr(packet, 'source_ip', '')[:15], 1)
                pdf.cell(35, 6, getattr(packet, 'destination_ip', '')[:15], 1)
                pdf.cell(20, 6, getattr(packet, 'protocol', ''), 1)
                pdf.cell(20, 6, str(getattr(packet, 'packet_size', 0)), 1)
                pdf.cell(20, 6, 'Suspicious' if getattr(packet, 'is_suspicious', False) else 'Normal', 1)
                pdf.ln()
        
        # Host Activity Details
        if data['host'] and len(data['host']) > 0:
            pdf.add_page()
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(0, 10, f'Host Activity Details (Top {min(50, len(data["host"]))})', 0, 1, 'L')
            pdf.set_font('Arial', '', 9)
            
            # Table header
            pdf.set_font('Arial', 'B', 9)
            pdf.cell(35, 8, 'Process Name', 1)
            pdf.cell(20, 8, 'PID', 1)
            pdf.cell(35, 8, 'User', 1)
            pdf.cell(20, 8, 'CPU %', 1)
            pdf.cell(20, 8, 'Memory %', 1)
            pdf.cell(25, 8, 'Status', 1)
            pdf.ln()
            
            pdf.set_font('Arial', '', 8)
            for proc in data['host'][:50]:
                pdf.cell(35, 6, getattr(proc, 'process_name', '')[:30], 1)
                pdf.cell(20, 6, str(getattr(proc, 'process_id', 0)), 1)
                pdf.cell(35, 6, getattr(proc, 'user', '')[:20], 1)
                pdf.cell(20, 6, f"{getattr(proc, 'cpu_usage', 0):.1f}", 1)
                pdf.cell(20, 6, f"{getattr(proc, 'memory_usage', 0):.1f}", 1)
                pdf.cell(25, 6, 'Suspicious' if getattr(proc, 'is_suspicious', False) else 'Normal', 1)
                pdf.ln()
        
        # Alerts Details
        if data['alerts'] and len(data['alerts']) > 0:
            pdf.add_page()
            pdf.set_font('Arial', 'B', 14)
            pdf.cell(0, 10, f'Alert Details (Top {min(30, len(data["alerts"]))})', 0, 1, 'L')
            pdf.set_font('Arial', '', 9)
            
            # Table header
            pdf.set_font('Arial', 'B', 9)
            pdf.cell(30, 8, 'Time', 1)
            pdf.cell(25, 8, 'Severity', 1)
            pdf.cell(30, 8, 'Type', 1)
            pdf.cell(85, 8, 'Description', 1)
            pdf.ln()
            
            pdf.set_font('Arial', '', 8)
            for alert in data['alerts'][:30]:
                timestamp = getattr(alert, 'timestamp', datetime.now())
                timestamp_str = timestamp.strftime('%m-%d %H:%M') if timestamp else ''
                
                severity = getattr(alert, 'severity', 'low')
                severity_display = severity.upper()
                
                pdf.cell(30, 6, timestamp_str, 1)
                pdf.cell(25, 6, severity_display, 1)
                pdf.cell(30, 6, getattr(alert, 'alert_type', '')[:15], 1)
                pdf.cell(85, 6, getattr(alert, 'description', '')[:70], 1)
                pdf.ln()
        
        pdf.output(filepath)
        print(f"✅ PDF report saved to {filepath}")
        return filepath
    
    def generate_ppt_report(self, report_type, start_date, end_date, data_source='database'):
        """Generate PowerPoint report"""
        print(f"📊 Generating {report_type} PPT report from {start_date} to {end_date} (source: {data_source})")
        
        data = self.get_data(start_date, end_date, data_source)
        stats = self.get_statistics(data)
        
        filename = f"report_{report_type}_{start_date.strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(self.report_dir, filename)
        
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Riley Falcon Security Services"
        subtitle.text = f"IDPS Security Report\n{start_date.strftime('%Y-%m-%d')} to {end_date.strftime('%Y-%m-%d')}\nData Source: {data_source.upper()}\nGenerated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add logo to title slide
        if os.path.exists(self.logo_path):
            try:
                left = Inches(8)
                top = Inches(0.5)
                slide.shapes.add_picture(self.logo_path, left, top, height=Inches(1))
            except:
                pass
        
        # Summary slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Executive Summary"
        content = slide.placeholders[1]
        
        # Add host data to summary
        host_text = f"""
HOST MONITORING:
• Processes: {stats['host']['total_processes']:,}
• Suspicious: {stats['host']['suspicious_processes']:,}
• Avg CPU: {stats['host']['avg_cpu']:.1f}%
• Avg Memory: {stats['host']['avg_memory']:.1f}%
• Active Users: {stats['host']['unique_users']}"""
        
        content.text = f"""ALERTS OVERVIEW:
• Total Alerts: {stats['alerts']['total']:,}
  - Critical: {stats['alerts']['critical']:,}
  - High: {stats['alerts']['high']:,}
  - Medium: {stats['alerts']['medium']:,}
  - Low: {stats['alerts']['low']:,}
• Active: {stats['alerts']['active']:,} | Resolved: {stats['alerts']['resolved']:,}

NETWORK ANALYSIS:
• Total Packets: {stats['network']['total_packets']:,}
• Suspicious: {stats['network']['suspicious_packets']:,}
• TCP: {stats['network']['tcp']:,}
• UDP: {stats['network']['udp']:,}
• ICMP: {stats['network']['icmp']:,}
• HTTP: {stats['network']['http']:,}
• HTTPS: {stats['network']['https']:,}
• DNS: {stats['network']['dns']:,}
• Data Transferred: {stats['network']['total_bytes'] / (1024*1024):.2f} MB
{host_text}"""
        
        # Add detailed data slide if there are alerts
        if stats['alerts']['total'] > 0:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Alert Details"
            content = slide.placeholders[1]
            
            alert_text = "Recent Alerts:\n\n"
            for alert in data['alerts'][:20]:
                timestamp = getattr(alert, 'timestamp', datetime.now())
                timestamp_str = timestamp.strftime('%Y-%m-%d %H:%M:%S') if timestamp else ''
                alert_text += f"• [{timestamp_str}] {getattr(alert, 'severity', 'unknown').upper()}: {getattr(alert, 'description', '')[:80]}\n"
            
            content.text = alert_text
        
        # Add host activity slide if there are processes
        if stats['host']['total_processes'] > 0:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Top Host Processes"
            content = slide.placeholders[1]
            
            process_text = "Top CPU-Consuming Processes:\n\n"
            for proc in data['host'][:20]:
                process_text += f"• {getattr(proc, 'process_name', 'Unknown')} (PID: {getattr(proc, 'process_id', 0)}) - CPU: {getattr(proc, 'cpu_usage', 0):.1f}% | Memory: {getattr(proc, 'memory_usage', 0):.1f}%"
                if getattr(proc, 'is_suspicious', False):
                    process_text += " [SUSPICIOUS]"
                process_text += "\n"
            
            content.text = process_text
        
        # Add network traffic slide if data exists
        if stats['network']['total_packets'] > 0:
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = "Network Traffic Summary"
            content = slide.placeholders[1]
            
            network_text = f"""Protocol Distribution:
• TCP: {stats['network']['tcp']:,} packets ({stats['network']['tcp']/max(stats['network']['total_packets'],1)*100:.1f}%)
• UDP: {stats['network']['udp']:,} packets ({stats['network']['udp']/max(stats['network']['total_packets'],1)*100:.1f}%)
• ICMP: {stats['network']['icmp']:,} packets ({stats['network']['icmp']/max(stats['network']['total_packets'],1)*100:.1f}%)

Service Distribution:
• HTTP: {stats['network']['http']:,} packets
• HTTPS: {stats['network']['https']:,} packets
• DNS: {stats['network']['dns']:,} packets

Top Source IPs:"""
            
            for i, (ip, count) in enumerate(list(stats['network']['top_sources'].items())[:5]):
                network_text += f"\n• {ip}: {count:,} packets"
            
            content.text = network_text
        
        prs.save(filepath)
        print(f"✅ PPT report saved to {filepath}")
        return filepath


# Alias for backward compatibility
ComprehensiveReportGenerator = ReportGenerator